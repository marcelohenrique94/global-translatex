import os
os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = 'API_KEY_HERE'

from google.cloud import translate_v2 as translate
from pptx import Presentation

# Inicializa o objeto de tradução
translate_client = translate.Client()

# Solicita ao usuário que especifique o caminho do arquivo PowerPoint a ser traduzido
powerpoint_path = input("Digite o caminho do arquivo PowerPoint a ser traduzido: ")

# Verifica se o arquivo PowerPoint existe
if not os.path.exists(powerpoint_path):
    print("Arquivo PowerPoint não encontrado.")
    exit()

# Carrega o arquivo PowerPoint em inglês
pr = Presentation(powerpoint_path)

# Cria um novo arquivo PowerPoint em português brasileiro
pr_br = Presentation()

# Loop através de cada slide e traduza o texto
for slide in pr.slides:
    new_slide = pr_br.slides.add_slide(slide.slide_layout)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                translated_text = translate_client.translate(run.text, target_language='pt-BR')['translatedText']
                # Correção aqui: usa o método `add_textbox()` em vez de `add_paragraph()`
                new_textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_textbox.text = translated_text

# Salva o novo arquivo PowerPoint em português brasileiro
output_path = os.path.splitext(powerpoint_path)[0] + '_pt-BR.pptx'
pr_br.save(output_path)

print("Arquivo PowerPoint traduzido salvo em", output_path)
